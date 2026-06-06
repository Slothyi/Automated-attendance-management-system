import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR

# ==========================================
# 🎨 COLOR PALETTE DESIGN SYSTEM
# ==========================================
BG_COLOR = RGBColor(15, 23, 42)          # Slate-900 (Dark background)
CARD_BG = RGBColor(30, 41, 59)           # Slate-800 (Card background)
ACCENT_CYAN = RGBColor(6, 182, 212)      # Cyan-500 (Primary Accent)
ACCENT_BLUE = RGBColor(59, 130, 246)     # Blue-500 (Secondary Accent)
TEXT_WHITE = RGBColor(248, 250, 252)     # Slate-50 (Primary Text)
TEXT_MUTED = RGBColor(148, 163, 184)     # Slate-400 (Secondary/Muted Text)
GREEN = RGBColor(34, 197, 94)            # Green-500 (Success indicators)
RED = RGBColor(239, 68, 68)              # Red-500 (Warnings/Limitations)

def set_dark_background(slide):
    """Fills the slide background with Slate-900 color."""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = BG_COLOR
    background.line.fill.background()  # Remove border
    return background

def add_clean_text(slide, text, left, top, width, height, font_size=12, font_color=TEXT_WHITE, bold=False, align=PP_ALIGN.LEFT):
    """Adds a single textbox with styled text to ensure robust word wrap and font styles."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = "Segoe UI"
    return txBox

def add_bullet_points(slide, bullets, left, top, width, height, font_size=12, font_color=TEXT_WHITE):
    """Adds a textbox filled with clean custom bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0.02)
    for i, b_text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(5)
        run = p.add_run()
        run.text = "•  " + b_text
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.name = "Segoe UI"
    return txBox

def draw_rounded_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=ACCENT_CYAN, border_width=1):
    """Draws a rounded rectangle representing a card, styled with a specific border and background."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    return shape

def draw_diagram_node(slide, left, top, width, height, text, bg_color=CARD_BG, border_color=ACCENT_CYAN, font_size=11, bold=True):
    """Draws a card shape and centers text inside it to represent a workflow/architecture diagram node."""
    shape = draw_rounded_card(slide, left, top, width, height, bg_color, border_color, border_width=1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = TEXT_WHITE
    run.font.bold = bold
    run.font.name = "Segoe UI"
    return shape

def draw_arrow_connector(slide, start_x, start_y, end_x, end_y, color=ACCENT_BLUE, line_width=1.5):
    """Draws a thin horizontal or vertical rectangle as a custom arrow line to avoid standard connector bugs."""
    if abs(start_x - end_x) > abs(start_y - end_y): # Horizontal line
        thickness = Inches(0.03)
        length = abs(start_x - end_x)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            min(start_x, end_x), start_y - (thickness / 2), length, thickness
        )
    else: # Vertical line
        thickness = Inches(0.03)
        length = abs(start_y - end_y)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            start_x - (thickness / 2), min(start_y, end_y), thickness, length
        )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def create_styled_slide(prs, title, category, slide_num):
    """Generates a base slide with a dark background, structured header, divider line, and footer."""
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Use blank layout
    set_dark_background(slide)
    
    # Header Category Label (small uppercase subtext)
    add_text_box = add_clean_text(
        slide, category.upper(), Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3),
        font_size=10, font_color=ACCENT_CYAN, bold=True
    )
    
    # Slide Title
    add_clean_text(
        slide, title, Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.6),
        font_size=26, font_color=TEXT_WHITE, bold=True
    )
    
    # Custom thin divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35), Inches(11.733), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_CYAN
    line.line.fill.background()
    
    # Slide Footer
    add_clean_text(
        slide, "AUTOMATED ATTENDANCE MANAGEMENT SYSTEM  |  TECHNICAL SHOWCASE", 
        Inches(0.8), Inches(7.0), Inches(8.0), Inches(0.3), font_size=8.5, font_color=TEXT_MUTED
    )
    add_clean_text(
        slide, f"{slide_num} / 29", 
        Inches(11.5), Inches(7.0), Inches(1.0), Inches(0.3), 
        font_size=10, font_color=ACCENT_CYAN, bold=True, align=PP_ALIGN.RIGHT
    )
    
    return slide


# ==========================================
# 🚀 INITIALIZE PRESENTATION
# ==========================================
prs = Presentation()
# Set widescreen (16:9) slide dimensions
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# ==========================================
# 🖥️ SLIDE 1: TITLE SLIDE (GLOWING GLASSMORPHIC ENERGETIC STYLE)
# ==========================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_background(slide1)

# Add large glowing decorative background circles
glow1 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(-1), Inches(6), Inches(6))
glow1.fill.solid()
glow1.fill.fore_color.rgb = RGBColor(15, 32, 67) # Deep Blue
glow1.line.fill.background()

glow2 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(3), Inches(5), Inches(5))
glow2.fill.solid()
glow2.fill.fore_color.rgb = RGBColor(12, 45, 60) # Deep Cyan
glow2.line.fill.background()

# Title layout
add_clean_text(slide1, "AUTOMATED ATTENDANCE\nMANAGEMENT SYSTEM", Inches(1.0), Inches(1.8), Inches(11.3), Inches(1.8), font_size=42, font_color=TEXT_WHITE, bold=True)
add_clean_text(slide1, "A Three-Factor Biometric, Geofenced, and BLE Proximity Verification Ecosystem", Inches(1.0), Inches(3.7), Inches(11.3), Inches(0.5), font_size=18, font_color=ACCENT_CYAN, bold=False)

# Accent line
line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(4.3), Inches(4.5), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_CYAN
line.line.fill.background()

# Meta details card
draw_rounded_card(slide1, Inches(1.0), Inches(4.8), Inches(11.333), Inches(1.6), bg_color=CARD_BG, border_color=ACCENT_BLUE)
add_clean_text(slide1, "DEVELOPED FOR:", Inches(1.3), Inches(5.0), Inches(4.0), Inches(0.3), font_size=10, font_color=ACCENT_CYAN, bold=True)
add_clean_text(slide1, "College Viva, Final Year Project, and Industry Demonstration", Inches(1.3), Inches(5.3), Inches(4.0), Inches(0.8), font_size=12, font_color=TEXT_WHITE)

add_clean_text(slide1, "CORE SYSTEM CAPABILITIES:", Inches(5.8), Inches(5.0), Inches(5.0), Inches(0.3), font_size=10, font_color=ACCENT_CYAN, bold=True)
bullets_s1 = [
    "GPS Geofencing with Native Android Anti-Mock Protection",
    "BLE Scan Matching Room Beacons & Visual Time-Sensitive OTPs",
    "On-device ML Kit Facial Liveness + Server-side Deep Face Recognition"
]
add_bullet_points(slide1, bullets_s1, Inches(5.8), Inches(5.3), Inches(6.2), Inches(1.0), font_size=10.5)


# ==========================================
# 🖥️ SLIDE 2: PROBLEM STATEMENT
# ==========================================
slide2 = create_styled_slide(prs, "The Attendance Dilemma in Higher Education", "Problem Statement", 2)
# Column 1: Context & Vulnerability
add_clean_text(slide2, "The Critical Cost of Manual Attendance", Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.4), font_size=18, font_color=ACCENT_CYAN, bold=True)
bullets_s2_left = [
    "Traditional manual roll calls take 10 to 15 minutes per lecture, wasting valuable instructional hours across semesters.",
    "Paper registers are fragile, prone to damage, and offer zero security controls against malicious modifications.",
    "Compiling reports manually creates immense administrative overhead, causing transcription errors and data delays.",
    "Physical tracking systems lack an audit trail, making it extremely difficult to verify records during inspections."
]
add_bullet_points(slide2, bullets_s2_left, Inches(0.8), Inches(2.4), Inches(5.5), Inches(4.0), font_size=12.5)

# Column 2: Danger cards
draw_rounded_card(slide2, Inches(6.8), Inches(1.8), Inches(5.7), Inches(1.5), bg_color=CARD_BG, border_color=RED)
add_clean_text(slide2, "Proxy Attendance (Buddy Punching)", Inches(7.0), Inches(1.9), Inches(5.3), Inches(0.3), font_size=14, font_color=RED, bold=True)
add_clean_text(slide2, "Absent students routinely ask classmates to mark them present on sign-in sheets or log responses using shared account links.", Inches(7.0), Inches(2.2), Inches(5.3), Inches(0.9), font_size=11, font_color=TEXT_MUTED)

draw_rounded_card(slide2, Inches(6.8), Inches(3.5), Inches(5.7), Inches(1.5), bg_color=CARD_BG, border_color=RED)
add_clean_text(slide2, "Location & Identity Fraud", Inches(7.0), Inches(3.6), Inches(5.3), Inches(0.3), font_size=14, font_color=RED, bold=True)
add_clean_text(slide2, "Students utilize GPS spoofing apps or display a photograph of an absent peer to simple camera scanners to falsify their presence.", Inches(7.0), Inches(3.9), Inches(5.3), Inches(0.9), font_size=11, font_color=TEXT_MUTED)

draw_rounded_card(slide2, Inches(6.8), Inches(5.2), Inches(5.7), Inches(1.5), bg_color=CARD_BG, border_color=RED)
add_clean_text(slide2, "QR Code Exploitation", Inches(7.0), Inches(5.3), Inches(5.3), Inches(0.3), font_size=14, font_color=RED, bold=True)
add_clean_text(slide2, "Static QR codes displayed on projectors are instantly photographed and shared via messaging groups, allowing remote sign-ins.", Inches(7.0), Inches(5.6), Inches(5.3), Inches(0.9), font_size=11, font_color=TEXT_MUTED)


# ==========================================
# 🖥️ SLIDE 3: EXISTING SYSTEM LIMITATIONS
# ==========================================
slide3 = create_styled_slide(prs, "Why Current Automated Systems Fail", "Existing System Limitations", 3)
add_clean_text(slide3, "Common Technical Flaws & Security Exploits", Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.4), font_size=16, font_color=TEXT_MUTED, bold=False)

# Card 1: Static QR Codes
draw_rounded_card(slide3, Inches(0.8), Inches(2.2), Inches(3.7), Inches(4.4), bg_color=CARD_BG, border_color=RED)
add_clean_text(slide3, "Static QR Scan Failures", Inches(1.0), Inches(2.4), Inches(3.3), Inches(0.4), font_size=16, font_color=RED, bold=True)
bullets_s3_1 = [
    "No Proximity Proof: QR images are sent to external student groups instantly.",
    "Absent students scan shared QR images from dorms or off-campus locations.",
    "Zero user validation occurs when the QR scanner acts independently of bios."
]
add_bullet_points(slide3, bullets_s3_1, Inches(1.0), Inches(3.0), Inches(3.3), Inches(3.4), font_size=12)

# Card 2: GPS Geofencing
draw_rounded_card(slide3, Inches(4.8), Inches(2.2), Inches(3.7), Inches(4.4), bg_color=CARD_BG, border_color=RED)
add_clean_text(slide3, "GPS Spoofing Flaws", Inches(5.0), Inches(2.4), Inches(3.3), Inches(0.4), font_size=16, font_color=RED, bold=True)
bullets_s3_2 = [
    "Android mock location services allow apps to inject custom latitude/longitude.",
    "GPS drift inside concrete buildings causes false-positives and rejects.",
    "No verification of whether the student is actually inside the room or outside."
]
add_bullet_points(slide3, bullets_s3_2, Inches(5.0), Inches(3.0), Inches(3.3), Inches(3.4), font_size=12)

# Card 3: Static Face Bio
draw_rounded_card(slide3, Inches(8.8), Inches(2.2), Inches(3.7), Inches(4.4), bg_color=CARD_BG, border_color=RED)
add_clean_text(slide3, "Static Biometric Flaws", Inches(9.0), Inches(2.4), Inches(3.3), Inches(0.4), font_size=16, font_color=RED, bold=True)
bullets_s3_3 = [
    "Facial scanners without liveness check can be fooled using prints/screens.",
    "Device sharing: one student logs into multiple accounts on their phone.",
    "Heavy server latency when matching raw images without local filtering."
]
add_bullet_points(slide3, bullets_s3_3, Inches(9.0), Inches(3.0), Inches(3.3), Inches(3.4), font_size=12)


# ==========================================
# 🖥️ SLIDE 4: PROPOSED SOLUTION
# ==========================================
slide4 = create_styled_slide(prs, "The Three-Factor Verification Shield", "Proposed Solution", 4)

# 4 horizontal cards representing the 4 shields
shields = [
    ("GPS GEOFENCING", "Checks device coordinates against campus boundaries using Android location services, checking and blocking mock providers natively.", ACCENT_BLUE),
    ("BLE ROOM PROXIMITY", "Scans BLE advertisements from the classroom beacon. Proves physical presence inside the lecture hall via signal attenuation.", ACCENT_CYAN),
    ("VISUAL DYNAMIC OTP", "A time-sensitive 5-digit hex code generated by the admin. Displayed to students in-class and expires within 10 minutes.", ACCENT_BLUE),
    ("LIVENESS & BIOMETRICS", "Google ML Kit prompts randomized head turns to verify liveness. Server performs deep 128D facial vector matching.", ACCENT_CYAN)
]

for idx, (title, desc, color) in enumerate(shields):
    left = Inches(0.8 + idx * 2.95)
    draw_rounded_card(slide4, left, Inches(1.8), Inches(2.8), Inches(4.7), bg_color=CARD_BG, border_color=color, border_width=1.5)
    
    # Bullet number
    add_clean_text(slide4, f"0{idx+1}", left + Inches(0.2), Inches(2.0), Inches(2.4), Inches(0.5), font_size=24, font_color=color, bold=True)
    # Title
    add_clean_text(slide4, title, left + Inches(0.2), Inches(2.7), Inches(2.4), Inches(0.6), font_size=15, font_color=TEXT_WHITE, bold=True)
    # Divider
    line = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.2), Inches(3.4), Inches(2.4), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    # Desc
    add_clean_text(slide4, desc, left + Inches(0.2), Inches(3.6), Inches(2.4), Inches(2.6), font_size=11.5, font_color=TEXT_MUTED)


# ==========================================
# 🖥️ SLIDE 5: PROJECT OVERVIEW
# ==========================================
slide5 = create_styled_slide(prs, "Project Overview & Ecosystem Mapping", "Project Overview", 5)

add_clean_text(slide5, "Ecosystem Integration & Workflow Loop", Inches(0.8), Inches(1.6), Inches(6.0), Inches(0.4), font_size=18, font_color=ACCENT_CYAN, bold=True)
bullets_s5 = [
    "Dynamic Classrooms: Faculty initializes secure, temporary attendance sessions that broadcast dynamic identifiers and visual OTPs.",
    "Sensor Fusion Client: Students utilize a custom Kotlin Android app integrating camera, GPS sensors, and Bluetooth hardware to scan environments.",
    "Asynchronous Processing: The FastAPI backend acts as a gateway, validating physical proximity, verifying identity vectors, and persisting logs.",
    "Anti-Proxy Validation: Verification is failed if any single credential check, location validation, or facial liveness scan is bypassed.",
    "Integrated Reporting: Automatic CSV logs are generated inside classes, tracking cumulative statistics, absences, and registration vectors."
]
add_bullet_points(slide5, bullets_s5, Inches(0.8), Inches(2.1), Inches(6.0), Inches(4.5), font_size=12)

# Right: Diagram showing basic loop
draw_rounded_card(slide5, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8), bg_color=CARD_BG, border_color=ACCENT_BLUE)
add_clean_text(slide5, "Ecosystem Workflow loop", Inches(7.5), Inches(2.0), Inches(4.7), Inches(0.3), font_size=14, font_color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)

# Nodes for loop
draw_diagram_node(slide5, Inches(8.9), Inches(2.5), Inches(2.0), Inches(0.6), "1. Faculty Starts Session", bg_color=BG_COLOR, border_color=ACCENT_CYAN)
draw_diagram_node(slide5, Inches(8.9), Inches(3.5), Inches(2.0), Inches(0.6), "2. Student Scans Environment", bg_color=BG_COLOR, border_color=ACCENT_CYAN)
draw_diagram_node(slide5, Inches(8.9), Inches(4.5), Inches(2.0), Inches(0.6), "3. ML Face & Liveness Engine", bg_color=BG_COLOR, border_color=ACCENT_BLUE)
draw_diagram_node(slide5, Inches(8.9), Inches(5.5), Inches(2.0), Inches(0.6), "4. Save DB & Append CSV", bg_color=BG_COLOR, border_color=GREEN)

# Connectors
draw_arrow_connector(slide5, Inches(9.9), Inches(3.1), Inches(9.9), Inches(3.5))
draw_arrow_connector(slide5, Inches(9.9), Inches(4.1), Inches(9.9), Inches(4.5))
draw_arrow_connector(slide5, Inches(9.9), Inches(5.1), Inches(9.9), Inches(5.5))


# ==========================================
# 🖥️ SLIDE 6: TECHNOLOGY STACK
# ==========================================
slide6 = create_styled_slide(prs, "Built with Modern SaaS Technologies", "Technology Stack", 6)

tech_col = [
    ("STUDENT MOBILE CLIENT", [
        "Kotlin: Primary language for Android Native App Development.",
        "CameraX Jetpack API: Robust and unified camera preview and image capture framework.",
        "Google ML Kit: Runs on-device face detection and head Euler angles extraction.",
        "BLE Scanning Engine: Bluetooth Low Energy scanner for beacon filtering.",
        "Retrofit & OkHttp: Secure network request management with dynamic headers."
    ], ACCENT_BLUE),
    ("API GATEWAY & BACKEND", [
        "FastAPI: Python framework for high-throughput, async request routing.",
        "Uvicorn: Lightning-fast ASGI web server implementation.",
        "JWT Authentication: Cryptographically signed security tokens for users.",
        "Passlib & BCrypt: Password hashing and security context storage.",
        "CSV Logger: File-append log engine representing physical sheets."
    ], ACCENT_CYAN),
    ("STORAGE & ML MATCHING", [
        "MongoDB NoSQL: JSON document database storing student and session records.",
        "dlib / Face Recognition: Computes 128D biometric face vectors on the server.",
        "OpenCV Python: Resizes and normalizes images for dlib performance.",
        "NumPy Library: Calculates fast vectorized Euclidean distance for matches.",
        "PyMongo: Low-latency driver connecting Python directly to database."
    ], ACCENT_BLUE)
]

for idx, (title, bullets, color) in enumerate(tech_col):
    left = Inches(0.8 + idx * 3.95)
    draw_rounded_card(slide6, left, Inches(1.8), Inches(3.8), Inches(4.8), bg_color=CARD_BG, border_color=color, border_width=1.5)
    add_clean_text(slide6, title, left + Inches(0.2), Inches(2.0), Inches(3.4), Inches(0.4), font_size=15, font_color=color, bold=True)
    
    # Horizontal line
    line = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.2), Inches(2.5), Inches(3.4), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    
    add_bullet_points(slide6, bullets, left + Inches(0.2), Inches(2.7), Inches(3.4), Inches(3.7), font_size=11)


# ==========================================
# 🖥️ SLIDE 7: SYSTEM ARCHITECTURE DIAGRAM
# ==========================================
slide7 = create_styled_slide(prs, "End-to-End System Architecture", "Technical Architecture", 7)

# Columns structure representation
# Draw 4 main layers
draw_rounded_card(slide7, Inches(0.8), Inches(1.8), Inches(2.6), Inches(4.8), bg_color=CARD_BG, border_color=ACCENT_BLUE)
add_clean_text(slide7, "Student Device", Inches(0.8), Inches(2.0), Inches(2.6), Inches(0.3), font_size=14, font_color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
draw_diagram_node(slide7, Inches(1.1), Inches(2.5), Inches(2.0), Inches(0.7), "Android App\n(Kotlin UI)", bg_color=BG_COLOR, border_color=ACCENT_CYAN)
draw_diagram_node(slide7, Inches(1.1), Inches(3.5), Inches(2.0), Inches(0.7), "ML Kit\n(Liveness Step)", bg_color=BG_COLOR, border_color=ACCENT_BLUE)
draw_diagram_node(slide7, Inches(1.1), Inches(4.5), Inches(2.0), Inches(0.7), "Location & BLE\n(GPS / BLE Scan)", bg_color=BG_COLOR, border_color=ACCENT_CYAN)
draw_diagram_node(slide7, Inches(1.1), Inches(5.5), Inches(2.0), Inches(0.7), "Retrofit Client\n(Network requests)", bg_color=BG_COLOR, border_color=ACCENT_BLUE)

draw_rounded_card(slide7, Inches(4.2), Inches(1.8), Inches(2.6), Inches(4.8), bg_color=CARD_BG, border_color=ACCENT_CYAN)
add_clean_text(slide7, "FastAPI Gateway", Inches(4.2), Inches(2.0), Inches(2.6), Inches(0.3), font_size=14, font_color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)
draw_diagram_node(slide7, Inches(4.5), Inches(2.5), Inches(2.0), Inches(0.7), "ASGI Web Server\n(Uvicorn)", bg_color=BG_COLOR, border_color=ACCENT_CYAN)
draw_diagram_node(slide7, Inches(4.5), Inches(3.5), Inches(2.0), Inches(0.7), "APIRouter\n(Route handling)", bg_color=BG_COLOR, border_color=ACCENT_BLUE)
draw_diagram_node(slide7, Inches(4.5), Inches(4.5), Inches(2.0), Inches(0.7), "Auth Middleware\n(JWT active_token)", bg_color=BG_COLOR, border_color=ACCENT_CYAN)
draw_diagram_node(slide7, Inches(4.5), Inches(5.5), Inches(2.0), Inches(0.7), "Controllers\n(Business logic)", bg_color=BG_COLOR, border_color=ACCENT_BLUE)

draw_rounded_card(slide7, Inches(7.6), Inches(1.8), Inches(2.6), Inches(4.8), bg_color=CARD_BG, border_color=ACCENT_BLUE)
add_clean_text(slide7, "Deep Learning Engine", Inches(7.6), Inches(2.0), Inches(2.6), Inches(0.3), font_size=13, font_color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
draw_diagram_node(slide7, Inches(7.9), Inches(2.8), Inches(2.0), Inches(0.7), "OpenCV Python\n(Image Normalizer)", bg_color=BG_COLOR, border_color=ACCENT_CYAN)
draw_diagram_node(slide7, Inches(7.9), Inches(3.8), Inches(2.0), Inches(0.7), "dlib Engine\n(HOG Face Detector)", bg_color=BG_COLOR, border_color=ACCENT_BLUE)
draw_diagram_node(slide7, Inches(7.9), Inches(4.8), Inches(2.0), Inches(0.7), "Face Encoder\n(128D vector extraction)", bg_color=BG_COLOR, border_color=ACCENT_CYAN)
draw_diagram_node(slide7, Inches(7.9), Inches(5.7), Inches(2.0), Inches(0.7), "NumPy Evaluator\n(Euclidean matching)", bg_color=BG_COLOR, border_color=ACCENT_BLUE)

draw_rounded_card(slide7, Inches(10.9), Inches(1.8), Inches(1.8), Inches(4.8), bg_color=CARD_BG, border_color=GREEN)
add_clean_text(slide7, "Storage Layer", Inches(10.9), Inches(2.0), Inches(1.8), Inches(0.3), font_size=14, font_color=GREEN, bold=True, align=PP_ALIGN.CENTER)
draw_diagram_node(slide7, Inches(11.1), Inches(2.8), Inches(1.4), Inches(0.9), "MongoDB\nAtlas DB", bg_color=BG_COLOR, border_color=GREEN)
draw_diagram_node(slide7, Inches(11.1), Inches(4.2), Inches(1.4), Inches(0.9), "Static File\nStore (JPG)", bg_color=BG_COLOR, border_color=ACCENT_CYAN)
draw_diagram_node(slide7, Inches(11.1), Inches(5.5), Inches(1.4), Inches(0.9), "CSV Log\nExporters", bg_color=BG_COLOR, border_color=ACCENT_BLUE)

# Connecting lines between columns
draw_arrow_connector(slide7, Inches(3.4), Inches(4.2), Inches(4.2), Inches(4.2), color=ACCENT_BLUE)
draw_arrow_connector(slide7, Inches(6.8), Inches(4.2), Inches(7.6), Inches(4.2), color=ACCENT_CYAN)
draw_arrow_connector(slide7, Inches(10.2), Inches(4.2), Inches(10.9), Inches(4.2), color=GREEN)


# ==========================================
# 🖥️ SLIDE 8: APPLICATION MODULES
# ==========================================
slide8 = create_styled_slide(prs, "Modular Application Design", "System Components", 8)

modules = [
    ("01", "ADMIN PORTAL & SYSTEM SETUP", [
        "Class Creation: Define subject names, branches, and schedules.",
        "Student Registration: Add roll numbers and names to class lists.",
        "Session Controller: Activate BLE codes and dynamic visual OTPs.",
        "Report Engine: View present/absent statistics and export CSVs."
    ]),
    ("02", "STUDENT DEVICE UTILITIES", [
        "Ecosystem Select: Log in as a student or connect as admin.",
        "GPS Verification: Pulls active GPS coordinates from Android.",
        "BLE Proximity Scan: Searches nearby BLE beacons for classroom validation.",
        "Visual OTP Verification: Input visual validation codes directly."
    ]),
    ("03", "COMPUTER VISION BIOMETRIC", [
        "ML Kit Analyzer: Local face search, checking yaw and pitch.",
        "Head-Turn Tracker: Evaluates random directions for anti-spoof.",
        "Face Cropper: Local scaling & cropping to minimize upload sizes.",
        "dlib Matcher: Server matches 128D face encodings within 0.5 limit."
    ]),
    ("04", "AUTHENTICATION & MIDDLEWARE", [
        "JWT Management: Secure token generation mapping to admin/student.",
        "Single-Device Lock: Server compares requests with DB active_token.",
        "Password Encryption: Bcrypt hashing applied during registration.",
        "CORS & Security Settings: Restricts API routing to trusted clients."
    ])
]

for idx, (num, title, bullets) in enumerate(modules):
    r = idx // 2
    c = idx % 2
    x = Inches(0.8 + c * 5.95)
    y = Inches(1.8 + r * 2.5)
    
    draw_rounded_card(slide8, x, y, Inches(5.7), Inches(2.3), bg_color=CARD_BG, border_color=ACCENT_CYAN)
    
    # Bullet Number
    add_clean_text(slide8, num, x + Inches(0.2), y + Inches(0.2), Inches(1.0), Inches(0.4), font_size=20, font_color=ACCENT_CYAN, bold=True)
    # Title
    add_clean_text(slide8, title, x + Inches(0.9), y + Inches(0.2), Inches(4.5), Inches(0.4), font_size=14, font_color=TEXT_WHITE, bold=True)
    # Divider
    line = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.2), y + Inches(0.7), Inches(5.3), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()
    # Bullets
    add_bullet_points(slide8, bullets, x + Inches(0.2), y + Inches(0.85), Inches(5.3), Inches(1.3), font_size=10.5)


# ==========================================
# 🖥️ SLIDE 9: DATABASE DESIGN (ER DIAGRAM)
# ==========================================
slide9 = create_styled_slide(prs, "MongoDB Collection Schemas & Relations", "Data Modeling", 9)

# We will draw a mock ER diagram with 6 boxes representing MongoDB collections
collections = [
    ("admins", "Primary Admin Details", ["_id: ObjectId", "name: String", "email: String (Unique)", "password: Hash (Bcrypt)"], Inches(0.8), Inches(1.8)),
    ("registered_students", "Onboarded Accounts", ["_id: ObjectId", "name / roll / email", "password: Hash", "face_encoding: Array[128]", "class_id: String", "active_token: JWT"], Inches(4.8), Inches(1.8)),
    ("students", "Authorized Rolls", ["_id: ObjectId", "name / roll", "class_id: String", "class_name: String", "attendance_status: Str"], Inches(8.8), Inches(1.8)),
    ("classes", "Class & Schedules", ["_id: ObjectId", "class_name / section", "department / semester", "created_by: AdminId", "students: Array[Objects]", "expires_at: DateTime"], Inches(0.8), Inches(4.4)),
    ("attendance", "Attendance Logs", ["_id: ObjectId", "student_id: ObjectId", "name / roll", "class_name: String", "date: String (YYYY-MM-DD)", "status: Str ('Present')"], Inches(4.8), Inches(4.4)),
    ("attendance_sessions", "Active Live Sessions", ["_id: ObjectId", "class_id: String", "session_code / session_uuid", "bluetooth_name: String", "classroom_beacon: String", "otp_code: String", "active: Boolean"], Inches(8.8), Inches(4.4))
]

for name, desc, fields, left, top in collections:
    draw_rounded_card(slide9, left, top, Inches(3.7), Inches(2.3), bg_color=CARD_BG, border_color=ACCENT_CYAN, border_width=1.5)
    add_clean_text(slide9, name.upper(), left + Inches(0.2), top + Inches(0.15), Inches(3.3), Inches(0.3), font_size=13, font_color=ACCENT_CYAN, bold=True)
    add_clean_text(slide9, desc, left + Inches(0.2), top + Inches(0.4), Inches(3.3), Inches(0.2), font_size=9.5, font_color=TEXT_MUTED, bold=False)
    
    # Line
    line = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.2), top + Inches(0.65), Inches(3.3), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()
    
    add_bullet_points(slide9, fields, left + Inches(0.2), top + Inches(0.75), Inches(3.3), Inches(1.4), font_size=9.5)


# ==========================================
# 🖥️ SLIDE 10: AUTHENTICATION FLOW
# ==========================================
slide10 = create_styled_slide(prs, "Secure Authentication & Single Device Locking", "Access Control", 10)

# Column 1: Verification steps
add_clean_text(slide10, "Student Onboarding (Face Enrollment)", Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.4), font_size=17, font_color=ACCENT_CYAN, bold=True)
bullets_s10_left = [
    "Pre-Authorization Check: Registration fails if the student's name and roll number are not pre-authorized by an admin in students collection.",
    "Initial Registration: Users input email, roll number, password, and upload a registration photo.",
    "Duplicate Face Check: Server extracts a 128D face vector. Scans database registered_students collection for any matching vector (Euclidean distance < 0.5). Blocks duplicate registrations.",
    "Bcrypt & JWT Setup: Password is hashed using bcrypt. Profile is stored in registered_students collection containing the facial vector."
]
add_bullet_points(slide10, bullets_s10_left, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.3), font_size=11.5)

# Column 2: Device Locking
add_clean_text(slide10, "Login & Session Device Locking (active_token)", Inches(6.8), Inches(1.7), Inches(5.7), Inches(0.4), font_size=17, font_color=ACCENT_BLUE, bold=True)
bullets_s10_right = [
    "Credential Check: User logs in with email and password via bcrypt comparison.",
    "Unique JWT Generation: FastAPI creates a unique JWT token embedding user ID and role ('student').",
    "Session Locking (Database Update): The database updates registered_students with the new JWT as the current 'active_token'.",
    "Authorization Guard (Middleware): Every request carries the token. The auth dependency compares the request token against active_token in the DB.",
    "Proxy Defense: If a student shares credentials, a login on a second device updates the active_token. The first device is immediately logged out."
]
add_bullet_points(slide10, bullets_s10_right, Inches(6.8), Inches(2.2), Inches(5.7), Inches(4.3), font_size=11.5)


# ==========================================
# 🖥️ SLIDE 11: ADMIN WORKFLOW
# ==========================================
slide11 = create_styled_slide(prs, "Faculty Admin Session Creation", "Process Design", 11)

add_clean_text(slide11, "Steps for starting an attendance session", Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.3), font_size=15, font_color=TEXT_MUTED)

# Flow diagram steps
steps_admin = [
    ("Step 1", "LOGIN", "Logs into Admin dashboard via secure bcrypt auth."),
    ("Step 2", "CREATE CLASS", "Defines name, section, year, and semester."),
    ("Step 3", "START SESSION", "Clicks 'Start' for a specific class to trigger session creation."),
    ("Step 4", "GENERATE META", "Server creates dynamic UUID, BLE identifier & 5-digit hex OTP."),
    ("Step 5", "BROADCAST/SHOW", "Beacon starts advertising BLE name. Visual OTP is shown on screen."),
    ("Step 6", "AUTO-EXPIRE", "The active session expires automatically in 10 minutes.")
]

for idx, (step_num, title, desc) in enumerate(steps_admin):
    left = Inches(0.8 + idx * 1.95)
    draw_rounded_card(slide11, left, Inches(2.4), Inches(1.8), Inches(3.8), bg_color=CARD_BG, border_color=ACCENT_CYAN, border_width=1.5)
    
    # Circle for Step number
    circle = slide11.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.4), Inches(2.6), Inches(1.0), Inches(1.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BG_COLOR
    circle.line.color.rgb = ACCENT_BLUE
    circle.line.width = Pt(1.5)
    # Circle text
    tf = circle.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = step_num
    run.font.size = Pt(12)
    run.font.color.rgb = ACCENT_CYAN
    run.font.bold = True
    
    add_clean_text(slide11, title, left + Inches(0.1), Inches(3.8), Inches(1.6), Inches(0.4), font_size=12, font_color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_clean_text(slide11, desc, left + Inches(0.15), Inches(4.3), Inches(1.5), Inches(1.7), font_size=10, font_color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    
    # Arrow connector to next card
    if idx < 5:
        draw_arrow_connector(slide11, left + Inches(1.8), Inches(4.3), left + Inches(1.95), Inches(4.3), color=ACCENT_BLUE)


# ==========================================
# 🖥️ SLIDE 12: STUDENT WORKFLOW
# ==========================================
slide12 = create_styled_slide(prs, "Student Attendance Marking Journey", "User Experience", 12)

add_clean_text(slide12, "Steps to log verification variables and verify attendance", Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.3), font_size=15, font_color=TEXT_MUTED)

steps_student = [
    ("Step 1", "JOIN CLASS", "Student opens Android app, selects class, and taps Mark Attendance."),
    ("Step 2", "GPS LOCK", "App pulls location. Coordinates must be within 200m of campus."),
    ("Step 3", "ENTER OTP", "Manually enters the 5-digit hex OTP shown by the professor."),
    ("Step 4", "BLE VERIFY", "App scans BLE advertisements. Finds matching room beacon."),
    ("Step 5", "LIVENESS SCAN", "Front camera opens. Prompts head turns to verify liveness."),
    ("Step 6", "MARK PRESENT", "Server processes face match, registers present status, logs in CSV.")
]

for idx, (step_num, title, desc) in enumerate(steps_student):
    left = Inches(0.8 + idx * 1.95)
    draw_rounded_card(slide12, left, Inches(2.4), Inches(1.8), Inches(3.8), bg_color=CARD_BG, border_color=ACCENT_BLUE, border_width=1.5)
    
    # Circle for Step number
    circle = slide12.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.4), Inches(2.6), Inches(1.0), Inches(1.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BG_COLOR
    circle.line.color.rgb = ACCENT_CYAN
    circle.line.width = Pt(1.5)
    
    # Circle text
    tf = circle.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = step_num
    run.font.size = Pt(12)
    run.font.color.rgb = ACCENT_BLUE
    run.font.bold = True
    
    add_clean_text(slide12, title, left + Inches(0.1), Inches(3.8), Inches(1.6), Inches(0.4), font_size=12, font_color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_clean_text(slide12, desc, left + Inches(0.15), Inches(4.3), Inches(1.5), Inches(1.7), font_size=10, font_color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    
    # Arrow connector to next card
    if idx < 5:
        draw_arrow_connector(slide12, left + Inches(1.8), Inches(4.3), left + Inches(1.95), Inches(4.3), color=ACCENT_CYAN)


# ==========================================
# 🖥️ SLIDE 13: BLUETOOTH ATTENDANCE WORKFLOW
# ==========================================
slide13 = create_styled_slide(prs, "Proximity Proof via BLE Beacons", "Communication Engineering", 13)

# Column 1: Details
add_clean_text(slide13, "Bluetooth Low Energy (BLE) Verification Core", Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.4), font_size=18, font_color=ACCENT_CYAN, bold=True)
bullets_s13 = [
    "Scan Window: Upon starting verification, the student app activates background BLE scanning for exactly 4 seconds.",
    "Name Filter: The app filters all nearby advertisements, checking for names matching the classroom beacon string (e.g. 'CLASS_CSE_A').",
    "Signal Strength (RSSI) Check: Device RSSI values are collected. High attenuation indicates long distance from the beacon.",
    "Proximity Proof: If the expected classroom beacon is detected in the scan map, it confirms the user is inside the physical room.",
    "Fail Handling: If the beacon is missing or RSSI drops below the room threshold, the scan fails. Blocks logins from adjacent corridors or hallways."
]
add_bullet_points(slide13, bullets_s13, Inches(0.8), Inches(2.2), Inches(5.8), Inches(4.5), font_size=12)

# Column 2: Diagram
draw_rounded_card(slide13, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8), bg_color=CARD_BG, border_color=ACCENT_BLUE)
add_clean_text(slide13, "BLE Verification Proximity Proof Flow", Inches(7.4), Inches(2.1), Inches(4.9), Inches(0.3), font_size=14, font_color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)

# Diagram shapes
draw_diagram_node(slide13, Inches(8.8), Inches(2.6), Inches(2.1), Inches(0.6), "Teacher Phone / Beacon\nCLASS_CSE_A Adv", bg_color=BG_COLOR, border_color=ACCENT_CYAN)
draw_diagram_node(slide13, Inches(8.8), Inches(3.6), Inches(2.1), Inches(0.6), "Student BLE Scanner\n4-Second Scan Map", bg_color=BG_COLOR, border_color=ACCENT_BLUE)
draw_diagram_node(slide13, Inches(8.8), Inches(4.6), Inches(2.1), Inches(0.6), "Name Filter &\nRSSI Distance Match", bg_color=BG_COLOR, border_color=ACCENT_CYAN)
draw_diagram_node(slide13, Inches(8.8), Inches(5.6), Inches(2.1), Inches(0.6), "Accept: Inside Room\nReject: Mismatch Error", bg_color=BG_COLOR, border_color=GREEN)

draw_arrow_connector(slide13, Inches(9.85), Inches(3.2), Inches(9.85), Inches(3.6), color=ACCENT_BLUE)
draw_arrow_connector(slide13, Inches(9.85), Inches(4.2), Inches(9.85), Inches(4.6), color=ACCENT_CYAN)
draw_arrow_connector(slide13, Inches(9.85), Inches(5.2), Inches(9.85), Inches(5.6), color=GREEN)


# ==========================================
# 🖥️ SLIDE 14: FACE RECOGNITION WORKFLOW
# ==========================================
slide14 = create_styled_slide(prs, "Biometric Face Recognition Flow", "Computer Vision & ML", 14)

add_clean_text(slide14, "Server-Side Facial Embedding Extraction & Distance Scoring", Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.4), font_size=18, font_color=ACCENT_CYAN, bold=True)
bullets_s14 = [
    "Image Upload: The client sends a cropped front face JPEG via a multipart form to the '/api/attendance/mark' route.",
    "Downscaling Preprocessing: The server resizes the image by 0.5x using OpenCV to speed up localization pipelines.",
    "HOG Face Locator: OpenCV and dlib utilize the Histogram of Oriented Gradients (HOG) algorithm to locate the face bounding box.",
    "128-Dimensional Vector: Dlib's deep learning face recognition model computes a 128D floating-point vector mapping unique facial nodes.",
    "Euclidean Distance Check: Calculates np.linalg.norm(stored_vector - upload_vector). Values < 0.5 indicate a match."
]
add_bullet_points(slide14, bullets_s14, Inches(0.8), Inches(2.2), Inches(5.8), Inches(4.5), font_size=12)

# Column 2: Diagram
draw_rounded_card(slide14, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8), bg_color=CARD_BG, border_color=ACCENT_BLUE)
add_clean_text(slide14, "Biometric Face Comparison Flow", Inches(7.4), Inches(2.0), Inches(4.9), Inches(0.3), font_size=14, font_color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)

draw_diagram_node(slide14, Inches(7.6), Inches(2.5), Inches(1.8), Inches(0.8), "Uploaded image\n(Multipart JPEG)", bg_color=BG_COLOR, border_color=ACCENT_CYAN)
draw_diagram_node(slide14, Inches(10.2), Inches(2.5), Inches(1.8), Inches(0.8), "Face Bounding Box\n(OpenCV / HOG)", bg_color=BG_COLOR, border_color=ACCENT_BLUE)
draw_diagram_node(slide14, Inches(8.9), Inches(3.6), Inches(2.0), Inches(0.8), "128D Face Embedding\n(dlib extraction)", bg_color=BG_COLOR, border_color=ACCENT_CYAN)
draw_diagram_node(slide14, Inches(8.9), Inches(4.7), Inches(2.0), Inches(0.8), "Euclidean Evaluator\nnp.linalg.norm()", bg_color=BG_COLOR, border_color=ACCENT_BLUE)
draw_diagram_node(slide14, Inches(8.9), Inches(5.8), Inches(2.0), Inches(0.6), "Distance < 0.5 Check", bg_color=BG_COLOR, border_color=GREEN)

draw_arrow_connector(slide14, Inches(9.4), Inches(2.9), Inches(10.2), Inches(2.9), color=ACCENT_BLUE)
draw_arrow_connector(slide14, Inches(9.9), Inches(3.3), Inches(9.9), Inches(3.6), color=ACCENT_CYAN)
draw_arrow_connector(slide14, Inches(9.9), Inches(4.4), Inches(9.9), Inches(4.7), color=ACCENT_BLUE)
draw_arrow_connector(slide14, Inches(9.9), Inches(5.5), Inches(9.9), Inches(5.8), color=GREEN)


# ==========================================
# 🖥️ SLIDE 15: ATTENDANCE VERIFICATION LOGIC
# ==========================================
slide15 = create_styled_slide(prs, "Multi-Factor Verification Funnel", "Ecosystem Security", 15)

# Left Column: Funnel Diagram (rectangles with decreasing width)
draw_rounded_card(slide15, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8), bg_color=CARD_BG, border_color=ACCENT_BLUE)
add_clean_text(slide15, "Verification Funnel Steps", Inches(1.0), Inches(2.0), Inches(5.1), Inches(0.3), font_size=15, font_color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)

# Draw Funnel levels
draw_diagram_node(slide15, Inches(1.2), Inches(2.5), Inches(4.7), Inches(0.5), "GPS Geofencing Check (<200 meters of college coords)", bg_color=BG_COLOR, border_color=ACCENT_CYAN, font_size=10.5)
draw_diagram_node(slide15, Inches(1.5), Inches(3.2), Inches(4.1), Inches(0.5), "Active Session Check & Classroom Beacon Match", bg_color=BG_COLOR, border_color=ACCENT_BLUE, font_size=10.5)
draw_diagram_node(slide15, Inches(1.8), Inches(3.9), Inches(3.5), Inches(0.5), "Visual dynamic OTP validation (10 min expiry)", bg_color=BG_COLOR, border_color=ACCENT_CYAN, font_size=10.5)
draw_diagram_node(slide15, Inches(2.1), Inches(4.6), Inches(2.9), Inches(0.5), "Android client head-turn Liveness Check", bg_color=BG_COLOR, border_color=ACCENT_BLUE, font_size=10.5)
draw_diagram_node(slide15, Inches(2.4), Inches(5.3), Inches(2.3), Inches(0.5), "Server Face vector matching (<0.5)", bg_color=BG_COLOR, border_color=GREEN, font_size=10.5)
draw_diagram_node(slide15, Inches(2.7), Inches(6.0), Inches(1.7), Inches(0.4), "Attendance Marked", bg_color=GREEN, border_color=TEXT_WHITE, font_size=10.5)

# Right Column: Explanations
add_clean_text(slide15, "Why the funnel structure is used", Inches(6.8), Inches(1.8), Inches(5.7), Inches(0.4), font_size=18, font_color=ACCENT_CYAN, bold=True)
bullets_s15_right = [
    "Efficient Cascade: Sensor validations are run hierarchically to filter out non-compliant requests before running heavy computer vision models.",
    "Resource Preservation: Server-side deep face recognition is only executed if the client has already passed GPS, BLE, OTP, and liveness steps.",
    "Cryptographic Integrity: All factors must compile correctly. Failure at any single point terminates the transaction, preventing false logs.",
    "Real-time Updates: When a student passes the funnel, their database state is changed to 'Present' and written to the class CSV."
]
add_bullet_points(slide15, bullets_s15_right, Inches(6.8), Inches(2.3), Inches(5.7), Inches(4.3), font_size=12)


# ==========================================
# 🖥️ SLIDE 16: SECURITY FEATURES
# ==========================================
slide16 = create_styled_slide(prs, "Multi-Layered Security Architecture", "System Security", 16)

sec_cards = [
    ("DATA SECURITY", [
        "Bcrypt Password Hash: User passwords are encrypted with dynamic salt cycles.",
        "JWT Tokens: Client requests contain cryptographically signed tokens.",
        "CORS Protocols: Backend restricts API cross-origin routing.",
        "Secure Storage: Face encodings are stored as numerical arrays, preventing image theft."
    ], ACCENT_BLUE),
    ("CLIENT SENSOR INTEGRITY", [
        "Anti-Mocking: Native Android code flags location.isFromMockProvider.",
        "BLE Proximity: Validates presence using signal power limits (RSSI).",
        "Visual OTP: Uses in-class displays to block remote code submissions.",
        "No Local Cache: Verification photos are deleted from memory after upload."
    ], ACCENT_CYAN),
    ("BIOMETRIC PROTECTION", [
        "Face Duplicate Filter: Blocks multiple profiles registering with same face.",
        "Euler Yaw Checks: Enforces randomized left/right head movements.",
        "dlib Distance Threshold: Evaluates vector margins at a strict 0.5 limit.",
        "Aspect Ratio Normalizer: Dynamic crop window limits spoofing borders."
    ], ACCENT_BLUE)
]

for idx, (title, bullets, color) in enumerate(sec_cards):
    left = Inches(0.8 + idx * 3.95)
    draw_rounded_card(slide16, left, Inches(1.8), Inches(3.8), Inches(4.8), bg_color=CARD_BG, border_color=color, border_width=1.5)
    add_clean_text(slide16, title, left + Inches(0.2), Inches(2.0), Inches(3.4), Inches(0.4), font_size=15, font_color=color, bold=True)
    
    # Line
    line = slide16.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.2), Inches(2.5), Inches(3.4), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    
    add_bullet_points(slide16, bullets, left + Inches(0.2), Inches(2.7), Inches(3.4), Inches(3.7), font_size=11)


# ==========================================
# 🖥️ SLIDE 17: ANTI-PROXY ATTENDANCE MECHANISM
# ==========================================
slide17 = create_styled_slide(prs, "Ecosystem Defenses Against Proxy Attendance", "Anti-Fraud Mechanisms", 17)

# Left Column: Vulnerability Explanations
add_clean_text(slide17, "Defeating Attendance Manipulation Hacks", Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.4), font_size=18, font_color=ACCENT_CYAN, bold=True)
bullets_s17 = [
    "GPS Mocking Bypass: Defeated! Using Android location API, requests containing mock coordinates are flagged and rejected immediately.",
    "QR Code Sharing: Defeated! QR scans are replaced by classroom BLE scans and visual OTPs, preventing remote links.",
    "Photo Spoofing: Defeated! On-device face detection requires active yaw turns (turning left/right) to prove liveness.",
    "Account Credentials Sharing: Defeated! During user login, the database records the unique active JWT token. Logins on new devices overwrite this token, invalidating old sessions."
]
add_bullet_points(slide17, bullets_s17, Inches(0.8), Inches(2.2), Inches(5.8), Inches(4.5), font_size=12)

# Right Column: Mitigation Matrix
draw_rounded_card(slide17, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8), bg_color=CARD_BG, border_color=ACCENT_BLUE)
add_clean_text(slide17, "Ecosystem Defeat Matrix", Inches(7.4), Inches(2.0), Inches(4.9), Inches(0.3), font_size=14, font_color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)

# Drawing mock table grid
table_y = [2.5, 3.0, 3.5, 4.0, 4.5, 5.0, 5.5, 6.0]
cols_x = [7.4, 9.4, 11.0]

# Table Header
draw_rounded_card(slide17, Inches(7.4), Inches(2.5), Inches(4.9), Inches(0.4), bg_color=BG_COLOR, border_color=ACCENT_BLUE)
add_clean_text(slide17, "Attack Vector", Inches(7.5), Inches(2.55), Inches(1.8), Inches(0.3), font_size=10.5, font_color=ACCENT_CYAN, bold=True)
add_clean_text(slide17, "Proposed Mitigation", Inches(9.5), Inches(2.55), Inches(1.5), Inches(0.3), font_size=10.5, font_color=ACCENT_CYAN, bold=True)
add_clean_text(slide17, "Status", Inches(11.1), Inches(2.55), Inches(1.1), Inches(0.3), font_size=10.5, font_color=GREEN, bold=True)

matrix_rows = [
    ("GPS spoof / mock locations", "Native Android flag check", "DEFEATED"),
    ("Sharing QR screen scans", "BLE room scan + visual OTP", "DEFEATED"),
    ("Printed peer photos", "Interactive liveness turns", "DEFEATED"),
    ("Sharing credentials", "Single active JWT session lock", "DEFEATED"),
    ("Multiple accounts", "Registration duplicate face block", "DEFEATED")
]

for idx, (attack, mit, status) in enumerate(matrix_rows):
    y = Inches(3.0 + idx * 0.7)
    draw_rounded_card(slide17, Inches(7.4), y, Inches(4.9), Inches(0.65), bg_color=BG_COLOR, border_color=ACCENT_BLUE)
    add_clean_text(slide17, attack, Inches(7.5), y + Inches(0.05), Inches(1.8), Inches(0.55), font_size=9.5, font_color=TEXT_WHITE, bold=True)
    add_clean_text(slide17, mit, Inches(9.5), y + Inches(0.05), Inches(1.5), Inches(0.55), font_size=9.5, font_color=TEXT_MUTED)
    add_clean_text(slide17, status, Inches(11.1), y + Inches(0.15), Inches(1.1), Inches(0.3), font_size=10, font_color=GREEN, bold=True)


# ==========================================
# 🖥️ SLIDE 18: PRIVACY, CONSENT, AND DATA PROTECTION (NEW)
# ==========================================
slide18 = create_styled_slide(prs, "Privacy, Consent, and Data Protection", "Data Ethics & Compliance", 18)

privacy_cards = [
    ("CONSENT & COLLECTION", [
        "Explicit Opt-In: Students must explicitly opt-in to face verification during account onboarding.",
        "Zero Image Persistence: Raw verification photos are processed in memory and immediately discarded.",
        "Alternative Verification: Manual override or OTP verification option is provided for users opting out."
    ], ACCENT_BLUE),
    ("DATA ENCRYPTION", [
        "Numerical Vectors: Stores only 128D mathematical vectors; raw images are never kept on the server.",
        "Encrypted Identifiers: Relational database links vectors to hashed profiles, securing student identities.",
        "Secure Transmissions: All client-server request payloads are encrypted via HTTPS protocols."
    ], ACCENT_CYAN),
    ("RETENTION & CONTROL", [
        "Retention Policy: Daily attendance logs are archived and purged from active DBs semester-wise.",
        "Self-Service Purge: Students can delete their biometric profiles, instantly removing vectors from NoSQL DB.",
        "Access Auditing: Administrative report access is logged to prevent internal data exposure."
    ], ACCENT_BLUE)
]

for idx, (title, bullets, color) in enumerate(privacy_cards):
    left = Inches(0.8 + idx * 3.95)
    draw_rounded_card(slide18, left, Inches(1.8), Inches(3.8), Inches(4.8), bg_color=CARD_BG, border_color=color, border_width=1.5)
    add_clean_text(slide18, title, left + Inches(0.2), Inches(2.0), Inches(3.4), Inches(0.4), font_size=15, font_color=color, bold=True)
    
    # Line
    line = slide18.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.2), Inches(2.5), Inches(3.4), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    
    add_bullet_points(slide18, bullets, left + Inches(0.2), Inches(2.7), Inches(3.4), Inches(3.7), font_size=11)


# ==========================================
# 🖥️ SLIDE 19: API ARCHITECTURE (WAS 18)
# ==========================================
slide19 = create_styled_slide(prs, "FastAPI REST API Architecture", "Backend Interfaces", 19)

api_groups = [
    ("AUTHENTICATION ROUTER (/api/auth)", [
        "POST /login: Authenticates student; returns JWT; updates active_token.",
        "POST /register: Saves account details, verifies rolls, saves face vector.",
        "JWT Verification Middleware: Intercepts requests to validate tokens."
    ]),
    ("ATTENDANCE CONTROLLER (/api/attendance)", [
        "POST /mark: Multi-factor verification gateway (GPS, BLE, OTP, face upload).",
        "POST /unmark: Allows admins to reset attendance status to Absent.",
        "GET /status: Returns daily presence status (Present/Absent).",
        "GET /history: Returns student's last 7 attendance logs."
    ]),
    ("SESSION MANAGEMENT (/api/session)", [
        "POST /start/{class_id}: Generates dynamic session identifiers and OTP codes.",
        "POST /stop/{class_id}: Deactivates attendance window, disabling code inputs.",
        "GET /active/{class_id}: Client pulls expected beacon and OTP code variables."
    ]),
    ("CLASS & STUDENT GROUP (/api/class)", [
        "POST /create: Initializes class templates mapping to departments.",
        "POST /add-students: Uploads roll lists; builds student collections.",
        "GET /all: Returns active classes mapped to the calling Admin's ID.",
        "GET /students/{class_id}: Returns detailed classroom lists and attendance stats."
    ])
]

for idx, (title, routes) in enumerate(api_groups):
    r = idx // 2
    c = idx % 2
    x = Inches(0.8 + c * 5.95)
    y = Inches(1.8 + r * 2.5)
    
    draw_rounded_card(slide19, x, y, Inches(5.7), Inches(2.3), bg_color=CARD_BG, border_color=ACCENT_CYAN)
    add_clean_text(slide19, title, x + Inches(0.2), y + Inches(0.15), Inches(5.3), Inches(0.4), font_size=13.5, font_color=ACCENT_CYAN, bold=True)
    
    # Line
    line = slide19.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.2), y + Inches(0.6), Inches(5.3), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()
    
    add_bullet_points(slide19, routes, x + Inches(0.2), y + Inches(0.7), Inches(5.3), Inches(1.5), font_size=10.5)


# ==========================================
# 🖥️ SLIDE 20: SYSTEM OPTIMIZATIONS (WAS 19)
# ==========================================
slide20 = create_styled_slide(prs, "High-Performance System Optimizations", "System Optimization", 20)

opts = [
    ("VECTORIZED FACE MATCH", "Face recognition utilizes NumPy-based Euclidean distance comparisons. Vector matching is run in under 15ms, minimizing CPU bottlenecks.", ACCENT_BLUE),
    ("IMAGE DOWNSCALING & PREPROCESSING", "Android crops photos to face oval boundaries before uploading. Server downscales images by 0.5x using OpenCV to speed up face localization.", ACCENT_CYAN),
    ("DATABASE INDEXES", "Indexes are created on roll and email collections in MongoDB, ensuring query times remain under 5ms even during concurrent student access.", ACCENT_BLUE),
    ("ASYNCHRONOUS EXPORT", "FastAPI writes student records to class CSV spreadsheets in the background, preventing writing operations from blocking API responses.", ACCENT_CYAN)
]

for idx, (title, desc, color) in enumerate(opts):
    r = idx // 2
    c = idx % 2
    x = Inches(0.8 + c * 5.95)
    y = Inches(1.8 + r * 2.5)
    
    draw_rounded_card(slide20, x, y, Inches(5.7), Inches(2.3), bg_color=CARD_BG, border_color=color, border_width=1.5)
    
    # Bullet number
    add_clean_text(slide20, f"0{idx+1}", x + Inches(0.25), y + Inches(0.2), Inches(1.0), Inches(0.4), font_size=18, font_color=color, bold=True)
    # Title
    add_clean_text(slide20, title, x + Inches(0.9), y + Inches(0.2), Inches(4.5), Inches(0.4), font_size=13.5, font_color=TEXT_WHITE, bold=True)
    # Line
    line = slide20.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.25), y + Inches(0.7), Inches(5.2), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    # Desc
    add_clean_text(slide20, desc, x + Inches(0.25), y + Inches(0.85), Inches(5.2), Inches(1.3), font_size=11, font_color=TEXT_MUTED)


# ==========================================
# 🖥️ SLIDE 21: DEPLOYMENT REQUIREMENTS (NEW)
# ==========================================
slide21 = create_styled_slide(prs, "Infrastructure & System Pre-requisites", "System Implementation", 21)

# Column 1: Server and Hosting
add_clean_text(slide21, "Server & Backend Infrastructure", Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.4), font_size=18, font_color=ACCENT_CYAN, bold=True)
bullets_s21_left = [
    "OS Environment: Ubuntu Server 22.04 LTS recommended, or any Docker-compatible virtual private server (VPS).",
    "Python ASGI Backend: Python 3.10+ installation with Uvicorn ASGI server for high-concurrency request routing.",
    "Database Deployment: MongoDB Community Edition for self-hosted setups, or MongoDB Atlas for managed cloud storage.",
    "Biometric Library Compilation: Pre-installed build tools (CMake, gcc/g++) required to compile C++ dlib dependencies."
]
add_bullet_points(slide21, bullets_s21_left, Inches(0.8), Inches(2.2), Inches(5.8), Inches(4.5), font_size=12)

# Column 2: Client and Hardware
add_clean_text(slide21, "Client Devices & In-Class Hardware", Inches(6.8), Inches(1.7), Inches(5.7), Inches(0.4), font_size=18, font_color=ACCENT_BLUE, bold=True)
bullets_s21_right = [
    "Android Client OS: Android 8.0 Oreo (API Level 26) or higher required to support CameraX API and BLE features.",
    "Hardware Permissions: User must grant Camera, Fine Location, and Bluetooth Scanning permissions during runtime.",
    "Classroom BLE Beacons: Bluetooth Low Energy beacons (iBeacon format) deployed and calibrated in each classroom.",
    "In-Class Displays: Projection screens or monitor setups to show the dynamic visual OTP code during verification windows."
]
add_bullet_points(slide21, bullets_s21_right, Inches(6.8), Inches(2.2), Inches(5.7), Inches(4.5), font_size=12)


# ==========================================
# 🖥️ SLIDE 22: DEPLOYMENT COST ANALYSIS (NEW)
# ==========================================
slide22 = create_styled_slide(prs, "Cost Structure & Financial Feasibility", "Financial Feasibility", 22)

# Left Column: Custom Cost Table
add_clean_text(slide22, "Estimated Operational Cost Breakdown", Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.4), font_size=18, font_color=ACCENT_CYAN, bold=True)

# Draw Table Background
draw_rounded_card(slide22, Inches(0.8), Inches(2.2), Inches(5.8), Inches(4.3), bg_color=CARD_BG, border_color=ACCENT_BLUE)

# Table Header
draw_rounded_card(slide22, Inches(0.9), Inches(2.4), Inches(5.6), Inches(0.4), bg_color=BG_COLOR, border_color=ACCENT_CYAN)
add_clean_text(slide22, "Component", Inches(1.0), Inches(2.45), Inches(1.8), Inches(0.3), font_size=11, font_color=TEXT_WHITE, bold=True)
add_clean_text(slide22, "Specification", Inches(2.9), Inches(2.45), Inches(2.2), Inches(0.3), font_size=11, font_color=TEXT_WHITE, bold=True)
add_clean_text(slide22, "Estimated Cost", Inches(5.2), Inches(2.45), Inches(1.2), Inches(0.3), font_size=11, font_color=GREEN, bold=True)

# Table Rows
cost_rows = [
    ("Cloud Server VM", "2 vCPU, 4GB RAM VPS", "$15 - $40 / mo"),
    ("NoSQL Database", "MongoDB Atlas (M10 tier)", "$0 - $30 / mo"),
    ("Physical BLE Beacons", "1 Unit per Classroom", "$15 / room (one-time)"),
    ("Domain & Security", "Custom Domain + SSL", "$10 / year"),
    ("User App Deployment", "Google Play Console fee", "$25 (one-time)")
]

for idx, (comp, spec, price) in enumerate(cost_rows):
    row_y = Inches(3.0 + idx * 0.6)
    draw_rounded_card(slide22, Inches(0.9), row_y, Inches(5.6), Inches(0.5), bg_color=BG_COLOR, border_color=ACCENT_BLUE)
    add_clean_text(slide22, comp, Inches(1.0), row_y + Inches(0.1), Inches(1.8), Inches(0.35), font_size=10, font_color=TEXT_WHITE, bold=True)
    add_clean_text(slide22, spec, Inches(2.9), row_y + Inches(0.1), Inches(2.2), Inches(0.35), font_size=9.5, font_color=TEXT_MUTED)
    add_clean_text(slide22, price, Inches(5.2), row_y + Inches(0.1), Inches(1.2), Inches(0.35), font_size=10, font_color=GREEN, bold=True)

# Right Column: ROI and Value
add_clean_text(slide22, "Economic Value & ROI Analysis", Inches(7.0), Inches(1.7), Inches(5.5), Inches(0.4), font_size=18, font_color=ACCENT_BLUE, bold=True)
bullets_s22_right = [
    "Zero License Fees: Built entirely on open-source dependencies (FastAPI, OpenCV, ML Kit), eliminating recurring user licensing costs.",
    "Hardware Cost Savings: Eliminates the need for dedicated biometric hardware terminals by leveraging students' existing mobile devices.",
    "Administrative Efficiency: Pays for itself rapidly by recovering over 15 hours of manual roll-call time per course per semester.",
    "Scale Flexibility: Free tier database and low-cost server tiers allow initial pilot testing without any upfront software capital expense."
]
add_bullet_points(slide22, bullets_s22_right, Inches(7.0), Inches(2.2), Inches(5.5), Inches(4.5), font_size=12)


# ==========================================
# 🖥️ SLIDE 23: TECHNICAL CHALLENGES & ENGINEERING SOLUTIONS (WAS 21)
# ==========================================
slide23 = create_styled_slide(prs, "Technical Challenges & Engineering Solutions", "Engineering Journey", 23)

challenges = [
    ("01", "Android BLE Scan Latency & Mismatch", 
     "Challenge: Android BLE scanner takes time to locate and return advertising beacons, causing timeouts.\nResolution: Switched scan settings to Low Latency Mode and created a 4-second scan accumulation map to filter room beacons."),
    ("02", "Camera Aspect-Ratio Image Warping", 
     "Challenge: Raw output from front cameras is stretched or distorted when loading face models, affecting accuracy.\nResolution: Programmed custom cropping logic mapping coordinates directly to CameraX viewfinders."),
    ("03", "GPS Drift in Concrete Buildings", 
     "Challenge: Thick concrete classroom walls block GPS signals, causing drift and false failures.\nResolution: Set GPS checks to a wider 200m range and combined location coordinates with local room BLE scans."),
    ("04", "Concurrency & Write-blocking", 
     "Challenge: Heavy database read/write requests slow down responses during large classroom attendance sessions.\nResolution: Setup indexed fields on student rolls and set FastAPI to write class CSVs asynchronously.")
]

for idx, (num, title, desc) in enumerate(challenges):
    r = idx // 2
    c = idx % 2
    x = Inches(0.8 + c * 5.95)
    y = Inches(1.8 + r * 2.5)
    
    draw_rounded_card(slide23, x, y, Inches(5.7), Inches(2.3), bg_color=CARD_BG, border_color=ACCENT_CYAN)
    
    # Number
    add_clean_text(slide23, num, x + Inches(0.2), y + Inches(0.2), Inches(1.0), Inches(0.4), font_size=20, font_color=ACCENT_CYAN, bold=True)
    # Title
    add_clean_text(slide23, title, x + Inches(0.9), y + Inches(0.2), Inches(4.5), Inches(0.4), font_size=13.5, font_color=TEXT_WHITE, bold=True)
    # Line
    line = slide23.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.2), y + Inches(0.7), Inches(5.3), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()
    # Desc
    add_clean_text(slide23, desc, x + Inches(0.2), y + Inches(0.85), Inches(5.3), Inches(1.3), font_size=11, font_color=TEXT_MUTED)


# ==========================================
# 🖥️ SLIDE 24: MEASURABLE BENEFITS & OPERATIONAL IMPACT (WAS 22)
# ==========================================
slide24 = create_styled_slide(prs, "Measurable Benefits & Operational Impact", "System Performance", 24)

# Column 1: Explanations
add_clean_text(slide24, "System Advantages & Value Delivery", Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.4), font_size=18, font_color=ACCENT_CYAN, bold=True)
bullets_s24 = [
    "Security Hardening: Eliminates buddy punching, QR sharing, and GPS spoofing completely.",
    "Administrative Time Savings: Reduces roll call times from 10 minutes to under 5 seconds per student.",
    "Paperless Tracking: Saves administrative hours by automating classroom statistics and monthly reports.",
    "Self-contained: Fully open-source dependencies run without expensive third-party APIs.",
    "Real-time visibility: Dashboard updates instantly when students pass verification."
]
add_bullet_points(slide24, bullets_s24, Inches(0.8), Inches(2.2), Inches(5.8), Inches(4.5), font_size=12)

# Column 2: 3 Big Metric Cards
metrics = [
    ("100%", "PROXY ELIMINATION", "Defeats all common attendance hacks successfully.", ACCENT_BLUE),
    ("< 5s", "VERIFICATION SPEED", "Fast and smooth student login verification.", ACCENT_CYAN),
    ("0", "MANUAL CALCULATION", "Automatic real-time Excel/CSV exports.", GREEN)
]

for idx, (val, title, desc, color) in enumerate(metrics):
    y = Inches(1.8 + idx * 1.6)
    draw_rounded_card(slide24, Inches(7.2), y, Inches(5.3), Inches(1.4), bg_color=CARD_BG, border_color=color, border_width=1.5)
    
    # Large Metric value
    add_clean_text(slide24, val, Inches(7.4), y + Inches(0.2), Inches(1.5), Inches(0.8), font_size=32, font_color=color, bold=True)
    # Title
    add_clean_text(slide24, title, Inches(9.1), y + Inches(0.25), Inches(3.2), Inches(0.3), font_size=13, font_color=TEXT_WHITE, bold=True)
    # Desc
    add_clean_text(slide24, desc, Inches(9.1), y + Inches(0.6), Inches(3.2), Inches(0.6), font_size=10.5, font_color=TEXT_MUTED)


# ==========================================
# 🖥️ SLIDE 25: NEXT-GEN ATTENDANCE ADVANCEMENTS (WAS 20)
# ==========================================
slide25 = create_styled_slide(prs, "Next-Gen Attendance Advancements", "Future Roadmap", 25)

future_cols = [
    ("EDGE FACE RECOGNITION", [
        "Run facial similarity checks locally using TensorFlow Lite on device.",
        "Saves bandwidth: Uploads only verification hashes instead of full photos.",
        "Zero server CV overhead: Scales easily to thousands of concurrent users."
    ], ACCENT_BLUE),
    ("MULTI-SPECTRAL LIVENESS", [
        "Incorporate IR/depth sensors on phones supporting advanced biometric API.",
        "Defeats high-res 3D masks and deepfake screens.",
        "Continuous verification: Track attendance status throughout the lecture."
    ], ACCENT_CYAN),
    ("AI PREDICTIVE ANALYTICS", [
        "Analyze student attendance trends using machine learning models.",
        "Predicts semester dropouts and alerts advisors automatically.",
        "Optimizes classroom schedules by identifying peaks in student absence."
    ], ACCENT_BLUE)
]

for idx, (title, bullets, color) in enumerate(future_cols):
    left = Inches(0.8 + idx * 3.95)
    draw_rounded_card(slide25, left, Inches(1.8), Inches(3.8), Inches(4.8), bg_color=CARD_BG, border_color=color, border_width=1.5)
    add_clean_text(slide25, title, left + Inches(0.2), Inches(2.0), Inches(3.4), Inches(0.4), font_size=15, font_color=color, bold=True)
    
    # Line
    line = slide25.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.2), Inches(2.5), Inches(3.4), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    
    add_bullet_points(slide25, bullets, left + Inches(0.2), Inches(2.7), Inches(3.4), Inches(3.7), font_size=11)


# ==========================================
# 🖥️ SLIDE 26: CURRENT SYSTEM LIMITATIONS (NEW)
# ==========================================
slide26 = create_styled_slide(prs, "Current System Limitations & Constraints", "Project Boundaries", 26)

limitations_cards = [
    ("BIOMETRIC & ENVIRONMENT", [
        "Lighting Dependency: Extreme low-light conditions in lecture halls can degrade facial recognition accuracy.",
        "Occlusion Sensitivity: Users must remove heavy facial coverings, sunglasses, or medical masks for biometric scanning.",
        "Device Alignment: Low-resolution front cameras may experience higher vector distance mismatch margins."
    ], RED),
    ("HARDWARE & NETWORKS", [
        "GPS Signal Blockage: Dense concrete classroom walls can block GPS signals, making BLE scans the critical fallback.",
        "Beacon Range Calibration: Requiring precise transmission power calibration to prevent corridor/adjacent room bleeding.",
        "Platform Lock-in: Native client software is currently restricted to Android devices, lacking iOS support."
    ], RED),
    ("OPERATIONAL SYSTEM", [
        "Onboarding Overhead: Requires a one-time onboarding enrollment step for face vector database initialization.",
        "Proxy Token Sharing: While JWT session locks block multi-device logins, they do not prevent students from sharing OTPs offsite if BLE is bypassed.",
        "Server Latency: Simultaneous uploads of image binaries at the end of a session can cause brief queue bottlenecks."
    ], RED)
]

for idx, (title, bullets, color) in enumerate(limitations_cards):
    left = Inches(0.8 + idx * 3.95)
    draw_rounded_card(slide26, left, Inches(1.8), Inches(3.8), Inches(4.8), bg_color=CARD_BG, border_color=color, border_width=1.5)
    add_clean_text(slide26, title, left + Inches(0.2), Inches(2.0), Inches(3.4), Inches(0.4), font_size=15, font_color=color, bold=True)
    
    # Line
    line = slide26.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.2), Inches(2.5), Inches(3.4), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    
    add_bullet_points(slide26, bullets, left + Inches(0.2), Inches(2.7), Inches(3.4), Inches(3.7), font_size=11)


# ==========================================
# 🖥️ SLIDE 27: CLIENT MOBILE INTERFACE WALKTHROUGH (WAS 23)
# ==========================================
slide27 = create_styled_slide(prs, "Client Mobile Interface Walkthrough", "UX Design Showcase", 27)

add_clean_text(slide27, "Key User Interface Screens & Interactive Elements", Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.3), font_size=15, font_color=TEXT_MUTED)

screens = [
    ("01", "SPLASH & ROLE SELECT", "Interactive role select with custom canvas mesh backgrounds."),
    ("02", "LOGIN & REGISTER", "Includes email validation, passwords, and face enrollment."),
    ("03", "STUDENT DASHBOARD", "Shows weekly graphs, active sessions, and attendance statuses."),
    ("04", "CODE INPUT DIALOG", "A Material Design popup checking dynamic visual OTP codes."),
    ("05", "BLE SCANNER VIEW", "Real-time background scanner filtering room beacons."),
    ("06", "LIVENESS OVERLAY", "Draws face framing ovals and arrow directions for head turns."),
    ("07", "SUCCESS SCREEN", "Lottie animation playing on verified attendance marks."),
    ("08", "REPORTS & EXPORT", "Class list overview displaying real-time CSV statistics.")
]

for idx, (num, title, desc) in enumerate(screens):
    r = idx // 4
    c = idx % 4
    x = Inches(0.8 + c * 2.95)
    y = Inches(2.2 + r * 2.3)
    
    draw_rounded_card(slide27, x, y, Inches(2.8), Inches(2.1), bg_color=CARD_BG, border_color=ACCENT_CYAN)
    # Number
    add_clean_text(slide27, num, x + Inches(0.2), y + Inches(0.15), Inches(1.0), Inches(0.3), font_size=16, font_color=ACCENT_CYAN, bold=True)
    # Title
    add_clean_text(slide27, title, x + Inches(0.2), y + Inches(0.55), Inches(2.4), Inches(0.5), font_size=11, font_color=TEXT_WHITE, bold=True)
    # Line
    line = slide27.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.2), y + Inches(1.15), Inches(2.4), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()
    # Desc
    add_clean_text(slide27, desc, x + Inches(0.2), y + Inches(1.25), Inches(2.4), Inches(0.75), font_size=9.5, font_color=TEXT_MUTED)


# ==========================================
# 🖥️ SLIDE 28: SECURING MODERN CLASSROOMS (WAS 24)
# ==========================================
slide28 = create_styled_slide(prs, "Securing Modern Classrooms", "Summary", 28)

# Left Column: Key Takeaways
add_clean_text(slide28, "Project Conclusion & Accomplishments", Inches(0.8), Inches(1.7), Inches(5.8), Inches(0.4), font_size=18, font_color=ACCENT_CYAN, bold=True)
bullets_s28 = [
    "Successfully designed, implemented, and verified a secure, multi-factor automated attendance management system.",
    "Integrated mobile sensors, local ML analytics, and server-side biometrics to create a robust security shield.",
    "Eliminated attendance vulnerabilities (GPS mocking, QR sharing, buddy punching, static photos).",
    "Greatly reduced administrative load for teachers by automating statistics and spreadsheet logs.",
    "Ecosystem ready for deployment across academic and enterprise training environments."
]
add_bullet_points(slide28, bullets_s28, Inches(0.8), Inches(2.2), Inches(5.8), Inches(4.5), font_size=12)

# Right Column: Big Final Card
draw_rounded_card(slide28, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8), bg_color=CARD_BG, border_color=GREEN, border_width=2)
add_clean_text(slide28, "KEY DELIVERABLE METRICS", Inches(7.4), Inches(2.2), Inches(4.9), Inches(0.4), font_size=18, font_color=GREEN, bold=True, align=PP_ALIGN.CENTER)

# List inside key card
takeaways = [
    "Secure: Low Euclidean matching limits.",
    "Fast: Complete verification takes under 5s.",
    "Automated: Logs compiled in real-time CSV formats.",
    "No Spoofing: Robust head-turning liveness detection.",
    "On-device Processing: Client-side facial analysis."
]
add_bullet_points(slide28, takeaways, Inches(7.8), Inches(2.8), Inches(4.2), Inches(3.2), font_size=13)


# ==========================================
# 🖥️ SLIDE 29: THANK YOU (WAS 25)
# ==========================================
slide29 = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_background(slide29)

# Glowing circles for presentation finish
glow3 = slide29.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.66), Inches(0.75), Inches(6), Inches(6))
glow3.fill.solid()
glow3.fill.fore_color.rgb = RGBColor(15, 35, 60)
glow3.line.fill.background()

add_clean_text(slide29, "THANK YOU!", Inches(1.0), Inches(2.2), Inches(11.3), Inches(0.8), font_size=44, font_color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_clean_text(slide29, "Questions & Feedback", Inches(1.0), Inches(3.1), Inches(11.3), Inches(0.4), font_size=18, font_color=ACCENT_CYAN, bold=False, align=PP_ALIGN.CENTER)

# Divider line
line = slide29.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.66), Inches(3.7), Inches(4.0), Inches(0.03))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_CYAN
line.line.fill.background()

# Final Contact Card
draw_rounded_card(slide29, Inches(3.16), Inches(4.2), Inches(7.0), Inches(1.8), bg_color=CARD_BG, border_color=ACCENT_BLUE)
add_clean_text(slide29, "AUTOMATED ATTENDANCE MANAGEMENT SYSTEM", Inches(3.36), Inches(4.4), Inches(6.6), Inches(0.3), font_size=13, font_color=ACCENT_CYAN, bold=True, align=PP_ALIGN.CENTER)

contact_bullets = [
    "Ecosystem Source: github.com/Slothyi/Automated-attendance-management-system",
    "Technical Architecture: Python FastAPI Gateway | Kotlin Android App",
    "Biometrics Engine: Google ML Kit Face Analysis & dlib 128D Matching"
]
add_bullet_points(slide29, contact_bullets, Inches(3.46), Inches(4.85), Inches(6.4), Inches(1.0), font_size=10.5)


# ==========================================
# 💾 SAVE PRESENTATION
# ==========================================
output_filename = "Automated_Attendance_Presentation.pptx"
prs.save(output_filename)
print(f"Presentation generated successfully as '{output_filename}'")
